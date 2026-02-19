"""
ClaudeBot - Enhanced Claude API service with tool use capabilities.
Deployed at /opt/claudebot on Ubuntu VPS, bound to 127.0.0.1:8020.
"""

import os
import json
import time
import hmac
import hashlib
import logging
import base64
import uuid
from io import BytesIO
from datetime import datetime, timezone
from functools import wraps
from collections import OrderedDict
from threading import Lock

from flask import Flask, request, jsonify
from dotenv import load_dotenv
import anthropic

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

load_dotenv()

CLAUDEBOT_SECRET = os.getenv("CLAUDEBOT_SECRET", "")
CLAUDE_API_KEY = os.getenv("CLAUDE_API_KEY", "")
SUPABASE_URL = os.getenv("SUPABASE_URL", "")
SUPABASE_KEY = os.getenv("SUPABASE_KEY", "")
PORT = int(os.getenv("PORT", "8020"))

CLAUDE_MODEL = "claude-sonnet-4-5-20250929"
MAX_TIMESTAMP_DRIFT = 30  # seconds
NONCE_CACHE_SIZE = 10000
OUTPUT_DIR = "/tmp/claudebot_output"

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger("claudebot")

# ---------------------------------------------------------------------------
# Flask app
# ---------------------------------------------------------------------------

app = Flask(__name__)

# ---------------------------------------------------------------------------
# Nonce cache (bounded LRU via OrderedDict)
# ---------------------------------------------------------------------------

_nonce_lock = Lock()
_nonce_cache: OrderedDict[str, float] = OrderedDict()


def _check_and_store_nonce(nonce: str) -> bool:
    """Return True if nonce is new (not replayed). Evicts oldest when full."""
    with _nonce_lock:
        if nonce in _nonce_cache:
            return False
        _nonce_cache[nonce] = time.time()
        while len(_nonce_cache) > NONCE_CACHE_SIZE:
            _nonce_cache.popitem(last=False)
        return True


# ---------------------------------------------------------------------------
# HMAC authentication middleware
# ---------------------------------------------------------------------------

def require_hmac(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        sig = request.headers.get("X-Api-Signature", "")
        ts = request.headers.get("X-Request-Timestamp", "")
        nonce = request.headers.get("X-Request-Nonce", "")

        if not sig or not ts or not nonce:
            log.warning("Auth failed: missing headers")
            return jsonify({"error": "Missing authentication headers (X-Api-Signature, X-Request-Timestamp, X-Request-Nonce)"}), 401

        # Timestamp freshness
        try:
            req_time = float(ts)
        except ValueError:
            return jsonify({"error": "X-Request-Timestamp must be a numeric Unix epoch"}), 401

        drift = abs(time.time() - req_time)
        if drift > MAX_TIMESTAMP_DRIFT:
            log.warning("Auth failed: timestamp drift %.1fs", drift)
            return jsonify({"error": f"Request timestamp too old ({drift:.0f}s drift, max {MAX_TIMESTAMP_DRIFT}s)"}), 401

        # Replay protection
        if not _check_and_store_nonce(nonce):
            log.warning("Auth failed: replayed nonce %s", nonce)
            return jsonify({"error": "Nonce already used (replay detected)"}), 401

        # HMAC verification
        body = request.get_data(as_text=True)
        message = f"{ts}:{nonce}:{body}"
        expected = hmac.new(
            CLAUDEBOT_SECRET.encode(),
            message.encode(),
            hashlib.sha256,
        ).hexdigest()

        if not hmac.compare_digest(sig, expected):
            log.warning("Auth failed: signature mismatch")
            return jsonify({"error": "Invalid HMAC signature"}), 401

        return f(*args, **kwargs)

    return decorated


# ---------------------------------------------------------------------------
# Tool definitions (Anthropic tool-use schema)
# ---------------------------------------------------------------------------

TOOLS = [
    {
        "name": "web_search",
        "description": "Search the web using DuckDuckGo and return top results.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query"},
            },
            "required": ["query"],
        },
    },
    {
        "name": "supabase_query",
        "description": "Run a query against a Supabase table. Supported operations: select, insert, update, delete.",
        "input_schema": {
            "type": "object",
            "properties": {
                "table": {"type": "string", "description": "Table name"},
                "operation": {
                    "type": "string",
                    "enum": ["select", "insert", "update", "delete"],
                    "description": "CRUD operation",
                },
                "data": {
                    "type": "object",
                    "description": "Payload — filters for select/delete, row data for insert, {match, values} for update.",
                },
            },
            "required": ["table", "operation"],
        },
    },
    {
        "name": "supabase_storage_read",
        "description": "Download a file from Supabase Storage and return its content as base64.",
        "input_schema": {
            "type": "object",
            "properties": {
                "bucket": {"type": "string"},
                "path": {"type": "string", "description": "Path inside the bucket"},
            },
            "required": ["bucket", "path"],
        },
    },
    {
        "name": "supabase_storage_write",
        "description": "Upload content to Supabase Storage.",
        "input_schema": {
            "type": "object",
            "properties": {
                "bucket": {"type": "string"},
                "path": {"type": "string"},
                "content": {"type": "string", "description": "Base64-encoded file content"},
            },
            "required": ["bucket", "path", "content"],
        },
    },
    {
        "name": "create_presentation",
        "description": "Generate a PPTX presentation and return the file path.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                        },
                        "required": ["title", "content"],
                    },
                    "description": "List of slide objects with title and content.",
                },
            },
            "required": ["title", "slides"],
        },
    },
    {
        "name": "create_document",
        "description": "Generate a DOCX document and return the file path.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "content": {"type": "string", "description": "Document body (plain text or markdown-ish)."},
            },
            "required": ["title", "content"],
        },
    },
    {
        "name": "create_spreadsheet",
        "description": "Generate an XLSX spreadsheet and return the file path.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "data": {
                    "type": "array",
                    "items": {
                        "type": "array",
                        "items": {"type": "string"},
                    },
                    "description": "2-D array: first row is headers, rest are data rows.",
                },
            },
            "required": ["title", "data"],
        },
    },
    {
        "name": "create_pdf",
        "description": "Generate a PDF document and return the file path.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "content": {"type": "string", "description": "Text content for the PDF."},
            },
            "required": ["title", "content"],
        },
    },
]

# ---------------------------------------------------------------------------
# Tool implementations
# ---------------------------------------------------------------------------

def _ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def tool_web_search(query: str) -> str:
    try:
        from duckduckgo_search import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5))
        if not results:
            return json.dumps({"results": [], "note": "No results found."})
        return json.dumps({"results": results})
    except Exception as exc:
        log.exception("web_search failed")
        return json.dumps({"error": str(exc)})


def _supa_client():
    from supabase import create_client
    return create_client(SUPABASE_URL, SUPABASE_KEY)


def tool_supabase_query(table: str, operation: str, data: dict | None = None) -> str:
    try:
        sb = _supa_client()
        data = data or {}

        if operation == "select":
            q = sb.table(table).select("*")
            for col, val in data.items():
                q = q.eq(col, val)
            res = q.execute()

        elif operation == "insert":
            res = sb.table(table).insert(data).execute()

        elif operation == "update":
            match = data.get("match", {})
            values = data.get("values", {})
            q = sb.table(table).update(values)
            for col, val in match.items():
                q = q.eq(col, val)
            res = q.execute()

        elif operation == "delete":
            q = sb.table(table).delete()
            for col, val in data.items():
                q = q.eq(col, val)
            res = q.execute()

        else:
            return json.dumps({"error": f"Unknown operation: {operation}"})

        return json.dumps({"data": res.data, "count": len(res.data)})
    except Exception as exc:
        log.exception("supabase_query failed")
        return json.dumps({"error": str(exc)})


def tool_supabase_storage_read(bucket: str, path: str) -> str:
    try:
        sb = _supa_client()
        data = sb.storage.from_(bucket).download(path)
        encoded = base64.b64encode(data).decode()
        return json.dumps({"content_base64": encoded, "size_bytes": len(data)})
    except Exception as exc:
        log.exception("supabase_storage_read failed")
        return json.dumps({"error": str(exc)})


def tool_supabase_storage_write(bucket: str, path: str, content: str) -> str:
    try:
        sb = _supa_client()
        raw = base64.b64decode(content)
        sb.storage.from_(bucket).upload(path, raw)
        return json.dumps({"success": True, "path": f"{bucket}/{path}", "size_bytes": len(raw)})
    except Exception as exc:
        log.exception("supabase_storage_write failed")
        return json.dumps({"error": str(exc)})


def tool_create_presentation(title: str, slides: list) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        _ensure_output_dir()
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        # Content slides
        layout = prs.slide_layouts[1]
        for s in slides:
            sl = prs.slides.add_slide(layout)
            sl.shapes.title.text = s.get("title", "")
            body = sl.placeholders[1]
            body.text = s.get("content", "")

        fname = f"{uuid.uuid4().hex[:8]}_{_safe_filename(title)}.pptx"
        fpath = os.path.join(OUTPUT_DIR, fname)
        prs.save(fpath)
        log.info("Created presentation: %s", fpath)
        return json.dumps({"file": fpath, "slides": len(slides) + 1})
    except Exception as exc:
        log.exception("create_presentation failed")
        return json.dumps({"error": str(exc)})


def tool_create_document(title: str, content: str) -> str:
    try:
        from docx import Document
        from docx.shared import Pt

        _ensure_output_dir()
        doc = Document()
        doc.add_heading(title, level=0)
        for para in content.split("\n"):
            if para.strip():
                doc.add_paragraph(para.strip())

        fname = f"{uuid.uuid4().hex[:8]}_{_safe_filename(title)}.docx"
        fpath = os.path.join(OUTPUT_DIR, fname)
        doc.save(fpath)
        log.info("Created document: %s", fpath)
        return json.dumps({"file": fpath})
    except Exception as exc:
        log.exception("create_document failed")
        return json.dumps({"error": str(exc)})


def tool_create_spreadsheet(title: str, data: list) -> str:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font

        _ensure_output_dir()
        wb = Workbook()
        ws = wb.active
        ws.title = title[:31]  # Excel sheet name limit

        for r_idx, row in enumerate(data, start=1):
            for c_idx, val in enumerate(row, start=1):
                cell = ws.cell(row=r_idx, column=c_idx, value=val)
                if r_idx == 1:
                    cell.font = Font(bold=True)

        fname = f"{uuid.uuid4().hex[:8]}_{_safe_filename(title)}.xlsx"
        fpath = os.path.join(OUTPUT_DIR, fname)
        wb.save(fpath)
        log.info("Created spreadsheet: %s", fpath)
        return json.dumps({"file": fpath, "rows": len(data)})
    except Exception as exc:
        log.exception("create_spreadsheet failed")
        return json.dumps({"error": str(exc)})


def tool_create_pdf(title: str, content: str) -> str:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet

        _ensure_output_dir()
        fname = f"{uuid.uuid4().hex[:8]}_{_safe_filename(title)}.pdf"
        fpath = os.path.join(OUTPUT_DIR, fname)

        doc = SimpleDocTemplate(fpath, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 0.3 * inch))

        for para in content.split("\n"):
            if para.strip():
                story.append(Paragraph(para.strip(), styles["BodyText"]))
                story.append(Spacer(1, 0.15 * inch))

        doc.build(story)
        log.info("Created PDF: %s", fpath)
        return json.dumps({"file": fpath})
    except Exception as exc:
        log.exception("create_pdf failed")
        return json.dumps({"error": str(exc)})


def _safe_filename(name: str) -> str:
    return "".join(c if c.isalnum() or c in "-_ " else "" for c in name).strip().replace(" ", "_")[:60]


# Dispatcher
TOOL_DISPATCH = {
    "web_search": lambda inp: tool_web_search(inp["query"]),
    "supabase_query": lambda inp: tool_supabase_query(inp["table"], inp["operation"], inp.get("data")),
    "supabase_storage_read": lambda inp: tool_supabase_storage_read(inp["bucket"], inp["path"]),
    "supabase_storage_write": lambda inp: tool_supabase_storage_write(inp["bucket"], inp["path"], inp["content"]),
    "create_presentation": lambda inp: tool_create_presentation(inp["title"], inp["slides"]),
    "create_document": lambda inp: tool_create_document(inp["title"], inp["content"]),
    "create_spreadsheet": lambda inp: tool_create_spreadsheet(inp["title"], inp["data"]),
    "create_pdf": lambda inp: tool_create_pdf(inp["title"], inp["content"]),
}


# ---------------------------------------------------------------------------
# Claude conversation loop with tool use
# ---------------------------------------------------------------------------

def chat_with_claude(user_message: str) -> dict:
    """Send a message to Claude and handle the full tool-use loop."""
    client = anthropic.Anthropic(api_key=CLAUDE_API_KEY)

    messages = [{"role": "user", "content": user_message}]
    system_prompt = (
        "You are ClaudeBot, an AI assistant with access to tools for web search, "
        "Supabase database operations, and document generation (PPTX, DOCX, XLSX, PDF). "
        "Use the available tools when the user's request requires them. "
        "Be concise, helpful, and strategic."
    )

    tool_calls_made = []
    max_iterations = 10

    for iteration in range(max_iterations):
        log.info("Claude API call iteration %d, messages=%d", iteration + 1, len(messages))

        response = client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=4096,
            system=system_prompt,
            tools=TOOLS,
            messages=messages,
        )

        # Collect text and tool-use blocks
        assistant_text = ""
        tool_use_blocks = []

        for block in response.content:
            if block.type == "text":
                assistant_text += block.text
            elif block.type == "tool_use":
                tool_use_blocks.append(block)

        # If stop_reason is "end_turn" or no tool calls, we're done
        if response.stop_reason == "end_turn" or not tool_use_blocks:
            return {
                "response": assistant_text,
                "tool_calls": tool_calls_made,
                "model": CLAUDE_MODEL,
            }

        # Append the full assistant message (with tool_use blocks) to conversation
        messages.append({"role": "assistant", "content": response.content})

        # Execute each tool and build tool_result blocks
        tool_results = []
        for tb in tool_use_blocks:
            tool_name = tb.name
            tool_input = tb.input
            log.info("Executing tool: %s", tool_name)

            handler = TOOL_DISPATCH.get(tool_name)
            if handler:
                result_str = handler(tool_input)
            else:
                result_str = json.dumps({"error": f"Unknown tool: {tool_name}"})

            tool_calls_made.append({"tool": tool_name, "input": tool_input, "result": result_str})
            tool_results.append({
                "type": "tool_result",
                "tool_use_id": tb.id,
                "content": result_str,
            })

        messages.append({"role": "user", "content": tool_results})

    # Exhausted iterations
    return {
        "response": assistant_text or "Reached maximum tool iterations.",
        "tool_calls": tool_calls_made,
        "model": CLAUDE_MODEL,
        "warning": "max_iterations_reached",
    }


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


@app.route("/chat", methods=["POST"])
@require_hmac
def chat():
    body = request.get_json(silent=True)
    if not body or "message" not in body:
        return jsonify({"error": "Request body must contain a 'message' field."}), 400

    user_message = body["message"].strip()
    if not user_message:
        return jsonify({"error": "Message cannot be empty."}), 400

    log.info("Chat request: %.100s...", user_message)

    try:
        result = chat_with_claude(user_message)
        return jsonify(result)
    except anthropic.AuthenticationError:
        log.error("Anthropic API key invalid")
        return jsonify({"error": "Claude API authentication failed. Check CLAUDE_API_KEY."}), 500
    except anthropic.RateLimitError:
        log.warning("Anthropic rate limit hit")
        return jsonify({"error": "Claude API rate limit reached. Try again shortly."}), 429
    except Exception as exc:
        log.exception("Chat endpoint error")
        return jsonify({"error": f"Internal error: {str(exc)}"}), 500


# ---------------------------------------------------------------------------
# Entrypoint
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    if not CLAUDEBOT_SECRET:
        log.warning("CLAUDEBOT_SECRET is not set — HMAC auth will reject all requests!")
    if not CLAUDE_API_KEY:
        log.warning("CLAUDE_API_KEY is not set — Claude API calls will fail!")

    log.info("Starting ClaudeBot on 127.0.0.1:%d", PORT)
    app.run(host="127.0.0.1", port=PORT, debug=False)
