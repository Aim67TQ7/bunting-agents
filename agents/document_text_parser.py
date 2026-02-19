"""
Document Text Parser Agent
Extracts and processes text from various document formats
"""

import os
import json
import mimetypes
from pathlib import Path
from typing import Dict, List, Optional, Any
import re

# Document parsing libraries
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import markdown
except ImportError:
    markdown = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None


class DocumentTextParser:
    """Agent for parsing text from various document formats"""
    
    def __init__(self):
        self.supported_formats = self._get_supported_formats()
        
    def _get_supported_formats(self) -> List[str]:
        """Get list of supported document formats based on available libraries"""
        formats = ['.txt', '.json', '.csv', '.log', '.md']
        
        if PyPDF2:
            formats.append('.pdf')
        if Document:
            formats.extend(['.docx', '.doc'])
        if openpyxl:
            formats.extend(['.xlsx', '.xls'])
        if Presentation:
            formats.extend(['.pptx', '.ppt'])
        if BeautifulSoup:
            formats.extend(['.html', '.htm', '.xml'])
            
        return formats
    
    def parse(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        Parse text from a document
        
        Args:
            file_path: Path to the document
            options: Optional parsing options
                - encoding: Text encoding (default: utf-8)
                - max_pages: Maximum pages to parse (for PDFs)
                - extract_metadata: Extract document metadata
                - clean_text: Apply text cleaning
                
        Returns:
            Dictionary containing:
                - text: Extracted text content
                - metadata: Document metadata
                - structure: Document structure information
                - errors: Any parsing errors
        """
        options = options or {}
        result = {
            'text': '',
            'metadata': {},
            'structure': {},
            'errors': []
        }
        
        try:
            path = Path(file_path)
            if not path.exists():
                result['errors'].append(f"File not found: {file_path}")
                return result
                
            file_ext = path.suffix.lower()
            
            if file_ext not in self.supported_formats:
                result['errors'].append(f"Unsupported format: {file_ext}")
                return result
            
            # Parse based on file type
            if file_ext == '.txt':
                result = self._parse_text(path, options)
            elif file_ext == '.pdf':
                result = self._parse_pdf(path, options)
            elif file_ext in ['.docx', '.doc']:
                result = self._parse_word(path, options)
            elif file_ext in ['.xlsx', '.xls']:
                result = self._parse_excel(path, options)
            elif file_ext in ['.pptx', '.ppt']:
                result = self._parse_powerpoint(path, options)
            elif file_ext == '.json':
                result = self._parse_json(path, options)
            elif file_ext == '.csv':
                result = self._parse_csv(path, options)
            elif file_ext in ['.html', '.htm']:
                result = self._parse_html(path, options)
            elif file_ext == '.xml':
                result = self._parse_xml(path, options)
            elif file_ext == '.md':
                result = self._parse_markdown(path, options)
            elif file_ext == '.log':
                result = self._parse_log(path, options)
                
            # Apply text cleaning if requested
            if options.get('clean_text', False):
                result['text'] = self._clean_text(result['text'])
                
        except Exception as e:
            result['errors'].append(f"Parsing error: {str(e)}")
            
        return result
    
    def _parse_text(self, path: Path, options: Dict) -> Dict:
        """Parse plain text file"""
        encoding = options.get('encoding', 'utf-8')
        result = {
            'text': '',
            'metadata': {},
            'structure': {},
            'errors': []
        }
        
        try:
            with open(path, 'r', encoding=encoding) as f:
                result['text'] = f.read()
            
            result['metadata'] = {
                'format': 'text',
                'encoding': encoding,
                'size': path.stat().st_size,
                'lines': len(result['text'].split('\n'))
            }
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_pdf(self, path: Path, options: Dict) -> Dict:
        """Parse PDF document"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {'pages': []},
            'errors': []
        }
        
        if not PyPDF2:
            result['errors'].append("PyPDF2 not installed")
            return result
            
        try:
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                
                # Extract metadata
                if options.get('extract_metadata', True):
                    info = reader.metadata
                    if info:
                        result['metadata'] = {
                            'title': str(info.get('/Title', '')),
                            'author': str(info.get('/Author', '')),
                            'subject': str(info.get('/Subject', '')),
                            'creator': str(info.get('/Creator', '')),
                            'pages': len(reader.pages)
                        }
                
                # Extract text from pages
                max_pages = options.get('max_pages', len(reader.pages))
                pages_to_read = min(max_pages, len(reader.pages))
                
                text_content = []
                for i in range(pages_to_read):
                    page = reader.pages[i]
                    page_text = page.extract_text()
                    text_content.append(page_text)
                    
                    result['structure']['pages'].append({
                        'page_number': i + 1,
                        'text_length': len(page_text)
                    })
                
                result['text'] = '\n\n'.join(text_content)
                
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_word(self, path: Path, options: Dict) -> Dict:
        """Parse Word document"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {'paragraphs': [], 'headings': []},
            'errors': []
        }
        
        if not Document:
            result['errors'].append("python-docx not installed")
            return result
            
        try:
            doc = Document(path)
            
            # Extract metadata
            if options.get('extract_metadata', True):
                props = doc.core_properties
                result['metadata'] = {
                    'title': props.title or '',
                    'author': props.author or '',
                    'created': str(props.created) if props.created else '',
                    'modified': str(props.modified) if props.modified else '',
                    'paragraphs': len(doc.paragraphs)
                }
            
            # Extract text and structure
            text_content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
                    
                    # Identify headings
                    if para.style.name.startswith('Heading'):
                        result['structure']['headings'].append({
                            'level': para.style.name,
                            'text': para.text[:100]
                        })
                    
                    result['structure']['paragraphs'].append({
                        'style': para.style.name,
                        'length': len(para.text)
                    })
            
            result['text'] = '\n\n'.join(text_content)
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_excel(self, path: Path, options: Dict) -> Dict:
        """Parse Excel spreadsheet"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {'sheets': []},
            'errors': []
        }
        
        if not openpyxl:
            result['errors'].append("openpyxl not installed")
            return result
            
        try:
            workbook = openpyxl.load_workbook(path, read_only=True)
            
            result['metadata'] = {
                'format': 'excel',
                'sheets': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames
            }
            
            text_content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"Sheet: {sheet_name}\n"
                
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join(str(cell) if cell else '' for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
                
                sheet_text += '\n'.join(rows)
                text_content.append(sheet_text)
                
                result['structure']['sheets'].append({
                    'name': sheet_name,
                    'rows': sheet.max_row,
                    'columns': sheet.max_column
                })
            
            result['text'] = '\n\n'.join(text_content)
            workbook.close()
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_powerpoint(self, path: Path, options: Dict) -> Dict:
        """Parse PowerPoint presentation"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {'slides': []},
            'errors': []
        }
        
        if not Presentation:
            result['errors'].append("python-pptx not installed")
            return result
            
        try:
            prs = Presentation(path)
            
            result['metadata'] = {
                'format': 'powerpoint',
                'slides': len(prs.slides)
            }
            
            text_content = []
            for i, slide in enumerate(prs.slides):
                slide_text = f"Slide {i + 1}:\n"
                
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text_runs.append(shape.text)
                
                slide_text += '\n'.join(text_runs)
                text_content.append(slide_text)
                
                result['structure']['slides'].append({
                    'slide_number': i + 1,
                    'shapes': len(slide.shapes),
                    'text_length': len(slide_text)
                })
            
            result['text'] = '\n\n'.join(text_content)
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_json(self, path: Path, options: Dict) -> Dict:
        """Parse JSON file"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {},
            'errors': []
        }
        
        try:
            with open(path, 'r', encoding=options.get('encoding', 'utf-8')) as f:
                data = json.load(f)
            
            result['text'] = json.dumps(data, indent=2)
            result['metadata'] = {
                'format': 'json',
                'keys': len(data) if isinstance(data, dict) else None,
                'items': len(data) if isinstance(data, list) else None
            }
            result['structure'] = self._analyze_json_structure(data)
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_csv(self, path: Path, options: Dict) -> Dict:
        """Parse CSV file"""
        import csv
        
        result = {
            'text': '',
            'metadata': {},
            'structure': {},
            'errors': []
        }
        
        try:
            with open(path, 'r', encoding=options.get('encoding', 'utf-8')) as f:
                reader = csv.reader(f)
                rows = list(reader)
            
            result['text'] = '\n'.join([','.join(row) for row in rows])
            result['metadata'] = {
                'format': 'csv',
                'rows': len(rows),
                'columns': len(rows[0]) if rows else 0
            }
            
            if rows:
                result['structure'] = {
                    'headers': rows[0] if rows else [],
                    'data_rows': len(rows) - 1 if len(rows) > 1 else 0
                }
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_html(self, path: Path, options: Dict) -> Dict:
        """Parse HTML file"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {},
            'errors': []
        }
        
        if not BeautifulSoup:
            result['errors'].append("beautifulsoup4 not installed")
            return result
            
        try:
            with open(path, 'r', encoding=options.get('encoding', 'utf-8')) as f:
                soup = BeautifulSoup(f, 'html.parser')
            
            # Extract text
            result['text'] = soup.get_text(separator='\n', strip=True)
            
            # Extract metadata
            result['metadata'] = {
                'format': 'html',
                'title': soup.title.string if soup.title else '',
                'meta_tags': len(soup.find_all('meta'))
            }
            
            # Extract structure
            result['structure'] = {
                'headings': [h.get_text()[:100] for h in soup.find_all(['h1', 'h2', 'h3'])],
                'links': len(soup.find_all('a')),
                'images': len(soup.find_all('img')),
                'tables': len(soup.find_all('table'))
            }
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_xml(self, path: Path, options: Dict) -> Dict:
        """Parse XML file"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {},
            'errors': []
        }
        
        if not BeautifulSoup:
            result['errors'].append("beautifulsoup4 not installed")
            return result
            
        try:
            with open(path, 'r', encoding=options.get('encoding', 'utf-8')) as f:
                soup = BeautifulSoup(f, 'xml')
            
            result['text'] = soup.get_text(separator='\n', strip=True)
            result['metadata'] = {
                'format': 'xml',
                'root_tag': soup.find().name if soup.find() else None
            }
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_markdown(self, path: Path, options: Dict) -> Dict:
        """Parse Markdown file"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {'headings': []},
            'errors': []
        }
        
        try:
            with open(path, 'r', encoding=options.get('encoding', 'utf-8')) as f:
                content = f.read()
            
            result['text'] = content
            
            # Extract headings
            heading_pattern = r'^(#{1,6})\s+(.+)$'
            for match in re.finditer(heading_pattern, content, re.MULTILINE):
                level = len(match.group(1))
                text = match.group(2)
                result['structure']['headings'].append({
                    'level': level,
                    'text': text[:100]
                })
            
            result['metadata'] = {
                'format': 'markdown',
                'lines': len(content.split('\n')),
                'headings': len(result['structure']['headings'])
            }
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _parse_log(self, path: Path, options: Dict) -> Dict:
        """Parse log file"""
        result = {
            'text': '',
            'metadata': {},
            'structure': {'patterns': {}},
            'errors': []
        }
        
        try:
            with open(path, 'r', encoding=options.get('encoding', 'utf-8')) as f:
                content = f.read()
            
            result['text'] = content
            
            # Analyze log patterns
            lines = content.split('\n')
            result['metadata'] = {
                'format': 'log',
                'lines': len(lines),
                'size': path.stat().st_size
            }
            
            # Count common log levels
            log_levels = ['ERROR', 'WARNING', 'INFO', 'DEBUG', 'CRITICAL']
            for level in log_levels:
                count = sum(1 for line in lines if level in line.upper())
                if count > 0:
                    result['structure']['patterns'][level] = count
            
        except Exception as e:
            result['errors'].append(str(e))
            
        return result
    
    def _analyze_json_structure(self, data: Any, max_depth: int = 3, current_depth: int = 0) -> Dict:
        """Analyze JSON structure"""
        if current_depth >= max_depth:
            return {'type': type(data).__name__}
        
        if isinstance(data, dict):
            return {
                'type': 'object',
                'keys': list(data.keys())[:10],
                'key_count': len(data)
            }
        elif isinstance(data, list):
            return {
                'type': 'array',
                'length': len(data),
                'item_types': list(set(type(item).__name__ for item in data[:10]))
            }
        else:
            return {'type': type(data).__name__}
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove non-printable characters
        text = ''.join(char for char in text if char.isprintable() or char in '\n\t')
        
        # Normalize line breaks
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text.strip()
    
    def extract_sections(self, text: str, section_markers: Optional[List[str]] = None) -> Dict[str, str]:
        """Extract sections from text based on markers"""
        section_markers = section_markers or ['Introduction', 'Abstract', 'Conclusion', 'References']
        sections = {}
        
        for marker in section_markers:
            pattern = rf'(?i)(?:^|\n)({re.escape(marker)}.*?)(?:\n|$)(.*?)(?=(?:^|\n)(?:' + '|'.join(re.escape(m) for m in section_markers) + r')|\Z)'
            match = re.search(pattern, text, re.DOTALL)
            if match:
                sections[marker] = match.group(2).strip()
        
        return sections
    
    def get_statistics(self, text: str) -> Dict[str, Any]:
        """Get text statistics"""
        words = text.split()
        sentences = re.split(r'[.!?]+', text)
        
        return {
            'characters': len(text),
            'words': len(words),
            'sentences': len([s for s in sentences if s.strip()]),
            'paragraphs': len(text.split('\n\n')),
            'average_word_length': sum(len(word) for word in words) / len(words) if words else 0,
            'unique_words': len(set(words))
        }


if __name__ == "__main__":
    # Example usage
    parser = DocumentTextParser()
    
    # Print supported formats
    print("Supported formats:", parser.supported_formats)
    
    # Example parsing
    import sys
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        result = parser.parse(file_path, {'extract_metadata': True, 'clean_text': True})
        
        print(f"\nParsing: {file_path}")
        print(f"Text length: {len(result['text'])} characters")
        print(f"Metadata: {result['metadata']}")
        if result['errors']:
            print(f"Errors: {result['errors']}")
        else:
            print(f"First 500 characters:\n{result['text'][:500]}...")
            
            # Get statistics
            stats = parser.get_statistics(result['text'])
            print(f"\nStatistics: {stats}")